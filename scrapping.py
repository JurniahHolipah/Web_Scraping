from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
from wordcloud import WordCloud
import mysql.connector
import pandas as pd

# Koneksi ke MySQL
print("Mulai Koneksi MySQL...")
db = mysql.connector.connect(
    host="localhost",
    user="root",
    password="", 
    database="data_scraping"
)
print("Koneksi MYSQL berhasil!")
cursor = db.cursor()

# Ambil data aset dari database
cursor.execute("SELECT nama_aset FROM data_aset")
rows = cursor.fetchall()

# Convert ke DataFrame untuk analisis
df = pd.DataFrame(rows, columns=["nama_aset"])

# Analisis: Jumlah total dataset
total_datasets = len(df)
print(f"Total dataset: {total_datasets}")

# Buat objek Presentasi
prs = Presentation()

# Slide 1: Judul Presentasi
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Proyek Scraping Data Aset"
subtitle.text = "Menganalisis dan Visualisasi Data dari Data.gov"
print("Slide 1 (Judul) berhasil dibuat.")

# Slide 2: Tujuan Proyek
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
title.text = "Tujuan Proyek"
content = slide_2.shapes.placeholders[1]
content.text = "1. Melakukan scraping data aset dari Data.gov\n" \
               "2. Menyimpan hasil ke MySQL\n" \
               "3. Melakukan analisis dan visualisasi data\n" \
               "4. Membuat presentasi otomatis"
print("Slide 2 (Tujuan Proyek) berhasil dibuat.")

# Slide 3: Word Cloud Visualization
slide_3 = prs.slides.add_slide(prs.slide_layouts[5])
title = slide_3.shapes.title
title.text = "Word Cloud: Nama Aset"
print("Slide 3 (Word Cloud) berhasil dibuat.")

# Buat Word Cloud
wordcloud = WordCloud(width=800, height=400, background_color='white').generate(" ".join(df['nama_aset']))

# Simpan Word Cloud sebagai gambar sementara
wordcloud_image = "wordcloud_image.png"
wordcloud.to_file(wordcloud_image)
print(f"Word Cloud disimpan sebagai {wordcloud_image}")

# Masukkan gambar Word Cloud ke slide
slide_3.shapes.add_picture(wordcloud_image, Inches(0.5), Inches(1.5), width=Inches(9))

# Slide 4: Bar Chart Visualization
slide_4 = prs.slides.add_slide(prs.slide_layouts[5])
title = slide_4.shapes.title
title.text = "Top 10 Dataset Berdasarkan Panjang Nama Aset"
print("Slide 4 (Bar Chart) berhasil dibuat.")

# Buat Bar Chart
top_10_datasets = df['nama_aset'].str.len().sort_values(ascending=False).head(10)
top_10_datasets.plot(kind='bar', color='skyblue', figsize=(10, 5))
plt.title('Top 10 Dataset Berdasarkan Panjang Nama Aset')
plt.xlabel('Dataset')
plt.ylabel('Panjang Nama')
plt.xticks(rotation=45, ha='right')

# Simpan Bar Chart sebagai gambar sementara
bar_chart_image = "bar_chart_image.png"
plt.savefig(bar_chart_image, bbox_inches='tight')
plt.close()
print(f"Bar Chart disimpan sebagai {bar_chart_image}")

# Masukkan gambar Bar Chart ke slide
slide_4.shapes.add_picture(bar_chart_image, Inches(0.5), Inches(1.5), width=Inches(9))

# Slide 5: Penutupan
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
title.text = "Penutupan"
content = slide_5.shapes.placeholders[1]
content.text = "Terima kasih atas perhatian Anda!\n\nProyek ini menyajikan insight dan visualisasi data aset dari Data.gov."
print("Slide 5 (Penutupan) berhasil dibuat.")

# Simpan PowerPoint
presentation_file = 'presentation.pptx'
prs.save(presentation_file)
print(f"Presentasi berhasil disimpan sebagai {presentation_file}")

# Tutup koneksi database
cursor.close()
db.close()

print("Proses selesai.")
